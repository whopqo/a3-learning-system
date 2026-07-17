"""
PPTX导出工具 —— 专业设计，深蓝主题，中文字体适配
"""
import os, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from config import PROJ_ROOT

INK = RGBColor(0x1A, 0x23, 0x32)
INK_LIGHT = RGBColor(0x2D, 0x3F, 0x56)
GOLD = RGBColor(0xBF, 0x9B, 0x4E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
IVORY = RGBColor(0xF8, 0xF5, 0xF0)
BLACK = RGBColor(0x2D, 0x29, 0x24)
MUTED = RGBColor(0x8B, 0x85, 0x78)
LIGHT_BG = RGBColor(0xF0, 0xEC, 0xE2)

FONT = 'Microsoft YaHei'
FONT_SERIF = 'Georgia'

def _set_font(run, name=FONT, size=None, bold=False, color=None):
    """正确设置中西文字体——用SubElement而非set属性"""
    from lxml import etree
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag, ns in [('a:latin', 'latin'), ('a:ea', 'ea'), ('a:cs', 'cs')]:
        existing = rPr.find(qn(tag))
        if existing is not None:
            rPr.remove(existing)
        el = etree.SubElement(rPr, qn(tag))
        el.set('typeface', name)
    if size: run.font.size = size
    run.font.bold = bold
    if color: run.font.color.rgb = color

def _add_header_bar(slide, title, prs_width):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs_width, Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = INK; bar.line.fill.background()
    tf = bar.text_frame; tf.margin_left = Inches(1); tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]; p.text = title
    _set_font(p.runs[0] if p.runs else p.add_run(), size=Pt(28), bold=True, color=WHITE)

def _add_footer(slide, topic):
    ft = slide.shapes.add_textbox(Inches(0.8), Inches(6.85), Inches(11.5), Inches(0.4))
    p = ft.text_frame.paragraphs[0]; p.text = f"EduSynth AI学习助手 · {topic}"
    _set_font(p.runs[0] if p.runs else p.add_run(), size=Pt(9), color=MUTED)
    p.alignment = PP_ALIGN.CENTER

def _add_insight_box(slide, text, top):
    box = slide.shapes.add_shape(5, Inches(1.2), top, Inches(11), Inches(0.8))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xFE, 0xF9, 0xF0)
    box.line.color.rgb = GOLD; box.line.width = Pt(1.5)
    tf = box.text_frame; tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.1)
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = f"💡 {text[:200]}"
    _set_font(p.runs[0] if p.runs else p.add_run(), size=Pt(14), color=INK_LIGHT)

def export_pptx(ppt_outline: str, topic: str = "课件", output_dir: str = None) -> str:
    if output_dir is None: output_dir = os.path.join(PROJ_ROOT, "output")
    os.makedirs(output_dir, exist_ok=True)

    prs = Presentation()
    w, h = Inches(13.333), Inches(7.5)
    prs.slide_width = w; prs.slide_height = h

    slides_data = _parse(ppt_outline)
    if not slides_data:
        slides_data = [{"title": topic, "bullets": [l.strip() for l in ppt_outline.split("\n") if l.strip() and len(l.strip()) > 3]}]

    for i, sd in enumerate(slides_data):
        title = sd.get("title", f"第{i+1}页")
        bullets = sd.get("bullets", [])

        if i == 0:
            sl = prs.slides.add_slide(prs.slide_layouts[6])
            sl.background.fill.solid(); sl.background.fill.fore_color.rgb = INK

            # Gold decorative line
            line = sl.shapes.add_shape(1, Inches(4), Inches(2.2), Inches(5.333), Pt(3))
            line.fill.solid(); line.fill.fore_color.rgb = GOLD; line.line.fill.background()

            # Main title
            tb = sl.shapes.add_textbox(Inches(1.5), Inches(2.6), Inches(10.3), Inches(1.6))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = title; p.alignment = PP_ALIGN.CENTER
            _set_font(p.runs[0] if p.runs else p.add_run(), size=Pt(48), bold=True, color=WHITE)

            # Subtitle
            if bullets:
                tb2 = sl.shapes.add_textbox(Inches(2.5), Inches(4.4), Inches(8.3), Inches(0.8))
                p2 = tb2.text_frame.paragraphs[0]; p2.text = bullets[0][:100]; p2.alignment = PP_ALIGN.CENTER
                _set_font(p2.runs[0] if p2.runs else p2.add_run(), size=Pt(18), color=RGBColor(0xCA, 0xCF, 0xDA))

            # Bottom gold stripe
            bot = sl.shapes.add_shape(1, Inches(0), Inches(7.0), w, Pt(4))
            bot.fill.solid(); bot.fill.fore_color.rgb = GOLD; bot.line.fill.background()

            # Footer
            ft = sl.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.4))
            p3 = ft.text_frame.paragraphs[0]; p3.text = f"EduSynth AI学习助手 · {topic}"; p3.alignment = PP_ALIGN.CENTER
            _set_font(p3.runs[0] if p3.runs else p3.add_run(), size=Pt(10), color=RGBColor(0x88, 0x84, 0x78))

        elif not bullets:
            sl = prs.slides.add_slide(prs.slide_layouts[6])
            sl.background.fill.solid(); sl.background.fill.fore_color.rgb = IVORY
            _add_header_bar(sl, title, w)
            tb = sl.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.3), Inches(4.5))
            tb.text_frame.word_wrap = True
            _add_footer(sl, topic)

        else:
            # Content slide — alternating backgrounds
            bg_color = IVORY if i % 2 == 0 else WHITE
            sl = prs.slides.add_slide(prs.slide_layouts[6])
            sl.background.fill.solid(); sl.background.fill.fore_color.rgb = bg_color
            _add_header_bar(sl, title, w)

            # Gold accent line under header
            accent = sl.shapes.add_shape(1, Inches(1), Inches(1.15), Inches(2), Pt(3))
            accent.fill.solid(); accent.fill.fore_color.rgb = GOLD; accent.line.fill.background()

            # Bullet content
            tb = sl.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(10.5), Inches(4.5))
            tf = tb.text_frame; tf.word_wrap = True

            # 动态字号—内容多自动缩小
            max_bullets = min(len(bullets), 9)
            base_size = Pt(20) if max_bullets <= 5 else (Pt(16) if max_bullets <= 7 else Pt(14))
            for j, bullet in enumerate(bullets[:max_bullets]):
                # 安全去掉编号前缀(不误删合法内容如"3D卷积")
                clean = re.sub(r'^[\s\-\*]*\d+[\.、\)\]\s]\s*', '', bullet).strip()
                if re.match(r'^[\-\*]\s', clean):
                    clean = clean[1:].strip()
                if not clean: continue
                if j == 0: p = tf.paragraphs[0]
                else: p = tf.add_paragraph()
                p.text = clean
                is_main = j < min(4, max_bullets // 2 + 1)
                _set_font(p.runs[0] if p.runs else p.add_run(), size=base_size if is_main else Pt(base_size.pt - 4), bold=is_main, color=BLACK)
                p.space_after = Pt(12)
                p.level = 0

            # 金句框：最后一个bullet长度>30就加，不依赖关键词
            if bullets and len(bullets[-1]) > 30:
                last_raw = re.sub(r'^[\s\-\*]*\d*[\.、\)\]\s]*', '', bullets[-1]).strip()
                if len(last_raw) > 15:
                    _add_insight_box(sl, last_raw, Inches(5.6))

            _add_footer(sl, topic)

    # Thank you slide
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = INK

    gold_line = sl.shapes.add_shape(1, Inches(4), Inches(2.6), Inches(5.3), Pt(3))
    gold_line.fill.solid(); gold_line.fill.fore_color.rgb = GOLD; gold_line.line.fill.background()

    tb = sl.shapes.add_textbox(Inches(1.5), Inches(3.0), Inches(10.3), Inches(1.2))
    p = tb.text_frame.paragraphs[0]; p.text = "谢谢观看"; p.alignment = PP_ALIGN.CENTER
    _set_font(p.runs[0] if p.runs else p.add_run(), size=Pt(52), bold=True, color=WHITE)

    tb2 = sl.shapes.add_textbox(Inches(2), Inches(4.4), Inches(9.3), Inches(0.6))
    p2 = tb2.text_frame.paragraphs[0]; p2.text = f"EduSynth AI学习助手 · {topic}"; p2.alignment = PP_ALIGN.CENTER
    _set_font(p2.runs[0] if p2.runs else p2.add_run(), size=Pt(18), color=RGBColor(0xCA, 0xCF, 0xDA))

    safe_topic = re.sub(r'[^\w一-鿿]', '_', topic)[:30]
    safe_topic = safe_topic.strip('_') or "课件"
    import time as _time
    filepath = os.path.join(output_dir, f"{safe_topic}_{int(_time.time())%10000}_课件.pptx")
    try:
        prs.save(filepath)
    except Exception as e:
        raise RuntimeError(f"PPTX保存失败: {e}") from e
    return filepath

def _parse(text: str) -> list:
    slides = []; cur = None
    for line in text.split("\n"):
        line = line.strip()
        if not line: continue
        if re.match(r'^#{1,4}\s', line):
            if cur: slides.append(cur)
            cur = {"title": re.sub(r'^#{1,4}\s*', '', line).strip(), "bullets": []}
        elif cur is not None:
            cur["bullets"].append(line)
    if cur: slides.append(cur)
    return slides
